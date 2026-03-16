from __future__ import annotations

import argparse
import contextlib
import json
import re
import textwrap
import time
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from typing import Iterable
from urllib.parse import urlparse

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor as PptColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
OUTPUT_DIR = DOCS_DIR / "review-pack"
SCREENSHOT_DIR = OUTPUT_DIR / "screenshots"
DIAGRAM_DIR = OUTPUT_DIR / "diagrams"
VIDEO_DIR = OUTPUT_DIR / "video-assets"

NAVY = (16, 42, 67)
NAVY_SOFT = (229, 238, 247)
TEAL = (40, 134, 145)
TEAL_SOFT = (224, 245, 247)
ACCENT = (220, 146, 80)
ACCENT_SOFT = (252, 240, 228)
SLATE = (62, 77, 93)
LIGHT_BG = (246, 249, 252)
WHITE = (255, 255, 255)

IMAGE_PATTERN = re.compile(r"!\[(.*?)\]\((.*?)\)")
HEADING_PATTERN = re.compile(r"^(#{1,4})\s+(.*)$")
ORDERED_PATTERN = re.compile(r"^\d+\.\s+(.*)$")
BULLET_PATTERN = re.compile(r"^-\s+(.*)$")
TABLE_SEPARATOR_PATTERN = re.compile(r"^\s*\|(?:\s*:?-+:?\s*\|)+\s*$")


@dataclass
class SlideSpec:
    number: int
    title: str
    bullets: list[str]
    notes: str
    duration_seconds: float
    image: Path | None = None
    image_caption: str | None = None
    subtitle: str | None = None


def ensure_output_dirs() -> None:
    for directory in (OUTPUT_DIR, SCREENSHOT_DIR, DIAGRAM_DIR, VIDEO_DIR):
        directory.mkdir(parents=True, exist_ok=True)

    stale_paths = [
        OUTPUT_DIR / "TaskBuddy-Demo.mp4",
        OUTPUT_DIR / "TaskBuddy-Demo-Video-Status.txt",
        OUTPUT_DIR / "TaskBuddy-Demo-Deck.pptx",
        OUTPUT_DIR / "TaskBuddy-Demo-Script.docx",
        OUTPUT_DIR / "TaskBuddy-Manual-Test-Plan.docx",
        OUTPUT_DIR / "TaskBuddy-Technical-Documentation.docx",
        OUTPUT_DIR / "TaskBuddy-User-Guide.docx",
        OUTPUT_DIR / "taskbuddy-review.db",
        OUTPUT_DIR / "taskbuddy-runtime.db",
    ]
    for path in stale_paths:
        if path.exists():
            path.unlink()

    for path in VIDEO_DIR.iterdir() if VIDEO_DIR.exists() else []:
        if path.is_file():
            path.unlink()


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path("C:/Windows/Fonts/seguibl.ttf") if bold else Path("C:/Windows/Fonts/segoeui.ttf"),
        Path("C:/Windows/Fonts/calibrib.ttf") if bold else Path("C:/Windows/Fonts/calibri.ttf"),
        Path("C:/Windows/Fonts/arialbd.ttf") if bold else Path("C:/Windows/Fonts/arial.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def strip_inline_markdown(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    return text.replace("\\|", "|").strip()


def add_field(paragraph, instruction: str) -> None:
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    run._r.append(begin)

    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = instruction
    run._r.append(instr)

    separate = OxmlElement("w:fldChar")
    separate.set(qn("w:fldCharType"), "separate")
    run._r.append(separate)

    text = OxmlElement("w:t")
    text.text = " "
    run._r.append(text)

    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.append(end)


def apply_docx_theme(document: Document, subtitle: str) -> None:
    styles = document.styles
    normal = styles["Normal"]
    normal.font.name = "Segoe UI"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor(*SLATE)

    for style_name, size in (("Heading 1", 19), ("Heading 2", 15), ("Heading 3", 12), ("Heading 4", 11)):
        style = styles[style_name]
        style.font.name = "Segoe UI"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor(*NAVY)

    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    header = section.header.paragraphs[0]
    header.text = f"TaskBuddy | {subtitle}"
    header.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if header.runs:
        header.runs[0].font.size = Pt(8.5)
        header.runs[0].font.color.rgb = RGBColor(*SLATE)

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run("Page ")
    add_field(footer, "PAGE")


def add_cover_page(document: Document, subtitle: str, summary_line: str) -> None:
    title_paragraph = document.add_paragraph()
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_paragraph.add_run("TaskBuddy")
    title_run.font.name = "Segoe UI"
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(*NAVY)

    subtitle_paragraph = document.add_paragraph()
    subtitle_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle_paragraph.add_run(subtitle)
    subtitle_run.font.name = "Segoe UI"
    subtitle_run.font.size = Pt(18)
    subtitle_run.font.bold = True
    subtitle_run.font.color.rgb = RGBColor(*TEAL)

    author_paragraph = document.add_paragraph()
    author_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    author_run = author_paragraph.add_run("Prepared by Tony Joseph Sebastian")
    author_run.font.name = "Segoe UI"
    author_run.font.size = Pt(12)
    author_run.font.color.rgb = RGBColor(*SLATE)

    summary_paragraph = document.add_paragraph()
    summary_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    summary_run = summary_paragraph.add_run(summary_line)
    summary_run.font.name = "Segoe UI"
    summary_run.font.size = Pt(11)
    summary_run.font.color.rgb = RGBColor(*SLATE)

    metadata_table = document.add_table(rows=3, cols=2)
    metadata_table.style = "Table Grid"
    metadata_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    rows = [
        ("Document", subtitle),
        ("Version", "1.0"),
        ("Generated", datetime.now(UTC).astimezone().strftime("%Y-%m-%d %H:%M %Z")),
    ]
    for row_index, (label, value) in enumerate(rows):
        metadata_table.rows[row_index].cells[0].text = label
        metadata_table.rows[row_index].cells[1].text = value

    document.add_page_break()


def split_table_row(line: str) -> list[str]:
    return [strip_inline_markdown(cell.strip()) for cell in line.strip().strip("|").split("|")]


def render_markdown(document: Document, lines: list[str], markdown_dir: Path) -> None:
    index = 0
    while index < len(lines):
        line = lines[index].rstrip()
        stripped = line.strip()
        if not stripped:
            index += 1
            continue

        if stripped.startswith("```"):
            code_lines: list[str] = []
            index += 1
            while index < len(lines) and not lines[index].strip().startswith("```"):
                code_lines.append(lines[index].rstrip())
                index += 1
            paragraph = document.add_paragraph()
            run = paragraph.add_run("\n".join(code_lines))
            run.font.name = "Courier New"
            run.font.size = Pt(9)
            index += 1
            continue

        heading_match = HEADING_PATTERN.match(stripped)
        if heading_match:
            level = min(len(heading_match.group(1)), 4)
            document.add_heading(strip_inline_markdown(heading_match.group(2)), level=level)
            index += 1
            continue

        image_match = IMAGE_PATTERN.match(stripped)
        if image_match:
            caption = strip_inline_markdown(image_match.group(1))
            image_path = (markdown_dir / image_match.group(2)).resolve()
            paragraph = document.add_paragraph()
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if image_path.exists():
                paragraph.add_run().add_picture(str(image_path), width=Inches(6.35))
            else:
                paragraph.add_run(f"[Missing image: {image_path.name}]")
            caption_paragraph = document.add_paragraph(caption)
            caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if caption_paragraph.runs:
                caption_paragraph.runs[0].italic = True
                caption_paragraph.runs[0].font.size = Pt(9)
                caption_paragraph.runs[0].font.color.rgb = RGBColor(*SLATE)
            index += 1
            continue

        if stripped.startswith("|") and index + 1 < len(lines) and TABLE_SEPARATOR_PATTERN.match(lines[index + 1].strip()):
            headers = split_table_row(stripped)
            rows: list[list[str]] = []
            index += 2
            while index < len(lines) and lines[index].strip().startswith("|"):
                rows.append(split_table_row(lines[index].strip()))
                index += 1
            table = document.add_table(rows=1, cols=len(headers))
            table.style = "Table Grid"
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            for cell, header in zip(table.rows[0].cells, headers):
                cell.text = header
            for row in rows:
                cells = table.add_row().cells
                for cell, value in zip(cells, row):
                    cell.text = value
            document.add_paragraph()
            continue

        bullet_match = BULLET_PATTERN.match(stripped)
        if bullet_match:
            paragraph = document.add_paragraph(style="List Bullet")
            paragraph.add_run(strip_inline_markdown(bullet_match.group(1)))
            index += 1
            continue

        ordered_match = ORDERED_PATTERN.match(stripped)
        if ordered_match:
            paragraph = document.add_paragraph(style="List Number")
            paragraph.add_run(strip_inline_markdown(ordered_match.group(1)))
            index += 1
            continue

        paragraph_lines = [stripped]
        index += 1
        while index < len(lines):
            next_stripped = lines[index].strip()
            if (
                not next_stripped
                or next_stripped.startswith("```")
                or HEADING_PATTERN.match(next_stripped)
                or IMAGE_PATTERN.match(next_stripped)
                or BULLET_PATTERN.match(next_stripped)
                or ORDERED_PATTERN.match(next_stripped)
                or (next_stripped.startswith("|") and index + 1 < len(lines) and TABLE_SEPARATOR_PATTERN.match(lines[index + 1].strip()))
            ):
                break
            paragraph_lines.append(next_stripped)
            index += 1

        paragraph = document.add_paragraph(strip_inline_markdown(" ".join(paragraph_lines)))
        paragraph.paragraph_format.space_after = Pt(6)


def generate_docx(markdown_path: Path, output_path: Path, subtitle: str, summary_line: str) -> None:
    document = Document()
    apply_docx_theme(document, subtitle)
    add_cover_page(document, subtitle, summary_line)
    lines = markdown_path.read_text(encoding="utf-8").splitlines()
    body_lines = lines[1:] if lines and lines[0].startswith("# ") else lines
    render_markdown(document, body_lines, markdown_path.parent)
    document.save(output_path)


def create_base_canvas(title: str, subtitle: str, size: tuple[int, int] = (1600, 900)) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    image = Image.new("RGB", size, LIGHT_BG)
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((36, 36, size[0] - 36, size[1] - 36), radius=34, fill=WHITE, outline=NAVY, width=4)
    draw.text((88, 82), title, font=load_font(40, bold=True), fill=NAVY)
    draw.text((88, 142), subtitle, font=load_font(20), fill=SLATE)
    return image, draw


def draw_box(
    draw: ImageDraw.ImageDraw,
    rect: tuple[int, int, int, int],
    title: str,
    fill: tuple[int, int, int],
    outline: tuple[int, int, int] = NAVY,
) -> None:
    draw.rounded_rectangle(rect, radius=20, fill=fill, outline=outline, width=3)
    font = load_font(22, bold=True)
    x1, y1, _, _ = rect
    y = y1 + 18
    for line in textwrap.wrap(title, width=20):
        draw.text((x1 + 18, y), line, font=font, fill=SLATE)
        y += 28


def draw_arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: tuple[int, int, int] = NAVY) -> None:
    draw.line((start, end), fill=color, width=5)
    x1, y1 = start
    x2, y2 = end
    if abs(x2 - x1) >= abs(y2 - y1):
        offset = -16 if x2 >= x1 else 16
        points = [(x2, y2), (x2 + offset, y2 - 10), (x2 + offset, y2 + 10)]
    else:
        offset = -16 if y2 >= y1 else 16
        points = [(x2, y2), (x2 - 10, y2 + offset), (x2 + 10, y2 + offset)]
    draw.polygon(points, fill=color)


def parse_junit_testsuite(path: Path) -> tuple[int, int, float]:
    root = ET.parse(path).getroot()
    tests = int(root.attrib.get("tests", 0))
    failures = int(root.attrib.get("failures", 0)) + int(root.attrib.get("errors", 0))
    time_taken = float(root.attrib.get("time", 0.0))
    return tests, failures, time_taken


def generate_report_summary_image() -> None:
    backend_tests, backend_failures, backend_time = parse_junit_testsuite(ROOT / "reports" / "backend-junit.xml")
    frontend_tests, frontend_failures, frontend_time = parse_junit_testsuite(ROOT / "reports" / "frontend-junit.xml")

    image, draw = create_base_canvas("Testing Evidence", "Automated coverage, exported reports, and manual validation assets.")
    metric_font = load_font(42, bold=True)
    body_font = load_font(24)

    cards = [
        ("Backend tests", f"{backend_tests}", f"{backend_failures} failures • {backend_time:.2f}s", NAVY_SOFT),
        ("Frontend tests", f"{frontend_tests}", f"{frontend_failures} failures • {frontend_time:.2f}s", TEAL_SOFT),
        ("Artifacts", "HTML + XLSX + DOCX", "Dashboard, spreadsheet, deck, and docs pack", ACCENT_SOFT),
    ]
    x = 96
    for title, metric, detail, fill in cards:
        rect = (x, 230, x + 420, 470)
        draw.rounded_rectangle(rect, radius=22, fill=fill, outline=NAVY, width=3)
        draw.text((x + 26, 255), title, font=body_font, fill=SLATE)
        draw.text((x + 26, 315), metric, font=metric_font, fill=NAVY)
        draw.text((x + 26, 395), detail, font=body_font, fill=SLATE)
        x += 455

    footer = textwrap.fill(
        "TaskBuddy ships with backend unit and integration tests, frontend interaction coverage, JUnit XML exports, an HTML dashboard, an Excel summary, and a dedicated manual test plan.",
        width=84,
    )
    draw.text((100, 560), footer, font=body_font, fill=SLATE)
    image.save(DIAGRAM_DIR / "test-evidence.png")


def generate_diagrams() -> None:
    architecture, draw = create_base_canvas(
        "TaskBuddy Architecture",
        "FastAPI serves the UI and API while LangGraph orchestrates deterministic tools.",
    )
    boxes = {
        "React UI": (84, 220, 324, 340),
        "FastAPI routes": (392, 220, 668, 340),
        "AgentController": (736, 146, 1028, 266),
        "LangGraph": (1088, 146, 1390, 266),
        "TaskInterpreter": (1088, 322, 1390, 442),
        "Tool registry": (1088, 500, 1390, 620),
        "TaskRepository": (736, 392, 1028, 512),
        "SQLite": (736, 620, 1028, 740),
    }
    fills = {
        "React UI": NAVY_SOFT,
        "FastAPI routes": ACCENT_SOFT,
        "AgentController": NAVY_SOFT,
        "LangGraph": TEAL_SOFT,
        "TaskInterpreter": NAVY_SOFT,
        "Tool registry": TEAL_SOFT,
        "TaskRepository": NAVY_SOFT,
        "SQLite": ACCENT_SOFT,
    }
    for label, rect in boxes.items():
        draw_box(draw, rect, label, fills[label])
    draw_arrow(draw, (324, 280), (392, 280))
    draw_arrow(draw, (668, 248), (736, 206))
    draw_arrow(draw, (1028, 206), (1088, 206))
    draw_arrow(draw, (1238, 266), (1238, 322))
    draw_arrow(draw, (1238, 442), (1238, 500))
    draw_arrow(draw, (882, 266), (882, 392))
    draw_arrow(draw, (882, 512), (882, 620))
    architecture.save(DIAGRAM_DIR / "architecture-overview.png")

    lifecycle, draw = create_base_canvas(
        "Task Execution Lifecycle",
        "Each saved task follows validation, planning, tool execution, persistence, and UI updates.",
    )
    steps = [
        ("Browser submits a task", (90, 280, 330, 410), NAVY_SOFT),
        ("FastAPI validates the request", (392, 280, 672, 410), ACCENT_SOFT),
        ("LangGraph plans the tool flow", (734, 280, 1028, 410), TEAL_SOFT),
        ("Tools execute and retry once", (1090, 280, 1400, 410), NAVY_SOFT),
        ("Repository saves the turn", (1088, 560, 1400, 690), ACCENT_SOFT),
        ("UI renders final output and trace", (734, 560, 1028, 690), TEAL_SOFT),
    ]
    for label, rect, fill in steps:
        draw_box(draw, rect, label, fill)
    draw_arrow(draw, (330, 345), (392, 345))
    draw_arrow(draw, (672, 345), (734, 345))
    draw_arrow(draw, (1028, 345), (1090, 345))
    draw_arrow(draw, (1235, 410), (1235, 560))
    draw_arrow(draw, (1088, 625), (1028, 625))
    lifecycle.save(DIAGRAM_DIR / "request-lifecycle.png")

    database, draw = create_base_canvas(
        "SQLite Persistence",
        "Users own threads, threads contain saved turns, and turns record execution steps.",
    )
    entities = {
        "users\n- id\n- username\n- password_hash\n- role\n- created_at": (84, 206, 418, 468),
        "threads\n- id\n- user_id\n- title\n- created_at\n- updated_at": (474, 206, 818, 468),
        "task_turns\n- id\n- thread_id\n- raw_input\n- sanitized_input\n- status\n- final_output\n- trace_id": (890, 152, 1314, 540),
        "execution_steps\n- id\n- turn_id\n- step_number\n- phase\n- tool_name\n- status\n- message": (890, 608, 1314, 888),
    }
    fills = [NAVY_SOFT, TEAL_SOFT, ACCENT_SOFT, NAVY_SOFT]
    for (label, rect), fill in zip(entities.items(), fills, strict=True):
        draw_box(draw, rect, label, fill)
    draw_arrow(draw, (418, 338), (474, 338))
    draw_arrow(draw, (818, 276), (890, 276))
    draw_arrow(draw, (1102, 540), (1102, 608))
    database.save(DIAGRAM_DIR / "database-schema.png")

    project_layout, draw = create_base_canvas(
        "Project Layout",
        "Key folders, major files, and how the documentation pack fits into the repository.",
    )
    layout_boxes = [
        ("backend/\nFastAPI app, routes,\nagent orchestration,\npersistence, tools", (84, 218, 398, 470), NAVY_SOFT),
        ("frontend/\nReact workspace,\nAPI client, types,\nUI tests", (444, 218, 750, 470), TEAL_SOFT),
        ("docs/\nUser guide, technical\ncontent, manual tests,\ndemo script sources", (796, 218, 1116, 470), ACCENT_SOFT),
        ("scripts/\nRun helpers, export\nreports, build the\ndocumentation pack", (1162, 218, 1480, 470), NAVY_SOFT),
        ("tests/\nUnit + integration\ncoverage and metadata", (264, 588, 600, 790), TEAL_SOFT),
        ("reports/\nJUnit XML, HTML\ndashboard, Excel\nsummary", (684, 588, 1018, 790), ACCENT_SOFT),
        ("docs/review-pack/\nGenerated DOCX,\nPPTX, screenshots,\ndiagrams, and video", (1102, 588, 1480, 790), NAVY_SOFT),
    ]
    for label, rect, fill in layout_boxes:
        draw_box(draw, rect, label, fill)
    project_layout.save(DIAGRAM_DIR / "project-layout.png")

    generate_report_summary_image()


def build_demo_payloads() -> dict[str, object]:
    admin_user = {
        "user_id": "admin-user",
        "username": "admin",
        "role": "admin",
    }

    text_turn = {
        "turn_id": "turn-text-1",
        "task_text": 'Convert "task buddy" to uppercase',
        "status": "completed",
        "final_output": "TASK BUDDY",
        "output_data": {"operation": "uppercase", "text": "TASK BUDDY"},
        "tools_used": ["TextProcessorTool"],
        "execution_steps": [
            {
                "step_number": 1,
                "phase": "validation",
                "status": "completed",
                "message": "Validated and sanitized the text request.",
                "tool_name": None,
                "payload": {"sanitized_text": 'Convert "task buddy" to uppercase'},
            },
            {
                "step_number": 2,
                "phase": "planning",
                "status": "completed",
                "message": "Planned TextProcessorTool for uppercase conversion.",
                "tool_name": None,
                "payload": {"tool": "TextProcessorTool"},
            },
            {
                "step_number": 3,
                "phase": "execution",
                "status": "completed",
                "message": "Converted the quoted text to uppercase.",
                "tool_name": "TextProcessorTool",
                "payload": {"operation": "uppercase"},
            },
        ],
        "timestamp": "2026-03-15T12:01:00Z",
        "trace_id": "trace-text-1",
    }

    multi_turn = {
        "turn_id": "turn-multi-1",
        "task_text": "What is the weather in Toronto and calculate 25 * 3",
        "status": "completed",
        "final_output": "1. Toronto: Cloudy, 8C, humidity 71%.\n2. 75.0",
        "output_data": {
            "results": [
                {"city": "Toronto", "condition": "Cloudy", "temperature_c": 8, "humidity_pct": 71},
                {"expression": "25 * 3", "result": 75.0},
            ]
        },
        "tools_used": ["WeatherMockTool", "CalculatorTool"],
        "execution_steps": [
            {
                "step_number": 1,
                "phase": "validation",
                "status": "completed",
                "message": "Validated a two-step request.",
                "tool_name": None,
                "payload": {"subtasks": 2},
            },
            {
                "step_number": 2,
                "phase": "planning",
                "status": "completed",
                "message": "Planned WeatherMockTool followed by CalculatorTool.",
                "tool_name": None,
                "payload": {"tools": ["WeatherMockTool", "CalculatorTool"]},
            },
            {
                "step_number": 3,
                "phase": "execution",
                "status": "completed",
                "message": "Fetched the Toronto weather summary.",
                "tool_name": "WeatherMockTool",
                "payload": {"city": "Toronto"},
            },
            {
                "step_number": 4,
                "phase": "execution",
                "status": "completed",
                "message": "Calculated 25 * 3.",
                "tool_name": "CalculatorTool",
                "payload": {"expression": "25 * 3"},
            },
        ],
        "timestamp": "2026-03-15T12:05:00Z",
        "trace_id": "trace-multi-1",
    }

    text_thread = {
        "thread_id": "thread-text",
        "title": 'Convert "task buddy" to uppercase',
        "created_at": "2026-03-15T12:00:00Z",
        "updated_at": "2026-03-15T12:01:00Z",
        "turns": [text_turn],
    }

    multi_thread = {
        "thread_id": "thread-multi",
        "title": "What is the weather in Toronto and calculate 25 * 3",
        "created_at": "2026-03-15T12:04:00Z",
        "updated_at": "2026-03-15T12:05:00Z",
        "turns": [multi_turn],
    }

    thread_summaries = [
        {
            "thread_id": text_thread["thread_id"],
            "title": text_thread["title"],
            "last_message_preview": text_turn["task_text"],
            "updated_at": text_thread["updated_at"],
        },
        {
            "thread_id": multi_thread["thread_id"],
            "title": multi_thread["title"],
            "last_message_preview": multi_turn["task_text"],
            "updated_at": multi_thread["updated_at"],
        },
    ]

    users = [
        {
            "user_id": "admin-user",
            "username": "admin",
            "role": "admin",
            "created_at": "2026-03-15T11:50:00Z",
        },
        {
            "user_id": "user-1",
            "username": "buddyuser",
            "role": "user",
            "created_at": "2026-03-15T11:55:00Z",
        },
    ]

    return {
        "admin_user": admin_user,
        "thread_summaries": thread_summaries,
        "text_thread": text_thread,
        "multi_thread": multi_thread,
        "users": users,
    }


def capture_screenshots(base_url: str) -> None:
    payloads = build_demo_payloads()
    text_thread = payloads["text_thread"]
    multi_thread = payloads["multi_thread"]
    thread_summaries = payloads["thread_summaries"]
    admin_user = payloads["admin_user"]
    users = payloads["users"]

    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        context = browser.new_context(viewport={"width": 1440, "height": 980})
        page = context.new_page()
        auth_state = {"authenticated": False}

        def fulfill_json(route, payload: object, status: int = 200) -> None:
            route.fulfill(status=status, content_type="application/json", body=json.dumps(payload))

        def handle_api(route) -> None:
            request = route.request
            parsed = urlparse(request.url)
            path = parsed.path

            if path == "/api/v1/auth/me":
                if auth_state["authenticated"]:
                    fulfill_json(route, admin_user)
                else:
                    fulfill_json(
                        route,
                        {
                            "error_code": "AUTH_REQUIRED",
                            "message": "Authentication is required.",
                            "trace_id": "trace-auth-required",
                            "details": {},
                        },
                        status=401,
                    )
                return

            if path == "/api/v1/auth/login":
                auth_state["authenticated"] = True
                fulfill_json(route, admin_user)
                return

            if path == "/api/v1/auth/logout":
                auth_state["authenticated"] = False
                fulfill_json(route, {"status": "ok"})
                return

            if not auth_state["authenticated"]:
                fulfill_json(
                    route,
                    {
                        "error_code": "AUTH_REQUIRED",
                        "message": "Authentication is required.",
                        "trace_id": "trace-auth-required",
                        "details": {},
                    },
                    status=401,
                )
                return

            if path == "/api/v1/threads":
                fulfill_json(route, thread_summaries)
                return

            if path == f"/api/v1/threads/{text_thread['thread_id']}":
                fulfill_json(route, text_thread)
                return

            if path == f"/api/v1/threads/{multi_thread['thread_id']}":
                fulfill_json(route, multi_thread)
                return

            if path == "/api/v1/admin/users":
                fulfill_json(route, users)
                return

            fulfill_json(
                route,
                {
                    "error_code": "NOT_FOUND",
                    "message": "Mock endpoint not found.",
                    "trace_id": "trace-mock-not-found",
                    "details": {"path": path},
                },
                status=404,
            )

        page.route(f"{base_url}/api/v1/**", handle_api)

        page.goto(base_url, wait_until="networkidle")
        page.screenshot(path=str(SCREENSHOT_DIR / "login-page.png"), full_page=True)

        page.get_by_label("Username").fill("admin")
        page.get_by_label("Password").fill("admin123")
        page.get_by_role("button", name="Sign in").click()
        page.get_by_placeholder("Search chat threads").wait_for()
        page.screenshot(path=str(SCREENSHOT_DIR / "workspace-home.png"), full_page=True)

        page.goto(f"{base_url}/threads/{text_thread['thread_id']}", wait_until="networkidle")
        page.locator(".assistant-card__output-value", has_text="TASK BUDDY").last.wait_for()
        page.screenshot(path=str(SCREENSHOT_DIR / "completed-text-task.png"), full_page=True)

        page.goto(f"{base_url}/threads/{multi_thread['thread_id']}", wait_until="networkidle")
        page.locator(".assistant-card__output-value", has_text="Toronto: Cloudy, 8C, humidity 71%.").last.wait_for()
        page.screenshot(path=str(SCREENSHOT_DIR / "multi-tool-trace.png"), full_page=True)

        page.goto(f"{base_url}/admin", wait_until="networkidle")
        page.get_by_text("Admin - User Management").wait_for()
        page.screenshot(path=str(SCREENSHOT_DIR / "admin-page.png"), full_page=True)

        page.goto(base_url, wait_until="networkidle")
        composer = page.get_by_placeholder("Ask TaskBuddy to run up to 2 supported subtasks.")
        composer.fill('Convert "a and b" to uppercase and weather in Toronto and calculate 2+2')
        page.get_by_text("Use up to 2 subtasks in a single request.").wait_for()
        page.screenshot(path=str(SCREENSHOT_DIR / "validation-example.png"), full_page=True)

        browser.close()


def load_demo_notes() -> dict[str, str]:
    lines = (DOCS_DIR / "demo-script.md").read_text(encoding="utf-8").splitlines()
    notes: dict[str, str] = {}
    current_key: str | None = None
    buffer: list[str] = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("## "):
            if current_key and buffer:
                notes[current_key] = " ".join(buffer).strip()
            current_key = stripped.replace("## Slide ", "", 1).split(" - ", 1)[1]
            buffer = []
            continue
        if current_key and stripped:
            buffer.append(stripped)

    if current_key and buffer:
        notes[current_key] = " ".join(buffer).strip()

    return notes


def slide_specs() -> list[SlideSpec]:
    notes = load_demo_notes()
    return [
        SlideSpec(
            1,
            "TaskBuddy",
            ["FastAPI-served React workspace", "Deterministic tools with LangGraph orchestration", "Five-minute business and technical walkthrough"],
            notes["Title and summary"],
            15,
            SCREENSHOT_DIR / "workspace-home.png",
            "Workspace home view",
            "Prepared by Tony Joseph Sebastian",
        ),
        SlideSpec(
            2,
            "What TaskBuddy does",
            ["Signs users in with a cookie session", "Lets each user create and revisit chat threads", "Returns final output first with a visible execution trace"],
            notes["What TaskBuddy does"],
            20,
            SCREENSHOT_DIR / "completed-text-task.png",
            "Completed text task",
        ),
        SlideSpec(
            3,
            "Roles and protected access",
            ["Bootstrap admin manages local users", "Standard users work only inside their own threads", "Thread APIs require authentication and admin APIs require the admin role"],
            notes["Roles and protected access"],
            20,
            SCREENSHOT_DIR / "admin-page.png",
            "Admin user management page",
        ),
        SlideSpec(
            4,
            "End-user workflow",
            ["Sign in, create a chat, and submit a supported prompt", "Watch the response card update with final output, tools, structured data, and trace", "Reopen history later from the thread sidebar"],
            notes["End-user workflow"],
            35,
            SCREENSHOT_DIR / "multi-tool-trace.png",
            "Multi-tool response and trace",
        ),
        SlideSpec(
            5,
            "Tool catalog and sample prompts",
            ["Text processing, calculator, mock weather, currency conversion, and transaction categorization", "Prompt routing is deterministic and documented in the technical guide", "Manual test cases are included for every tool family"],
            notes["Tool catalog and sample prompts"],
            35,
            SCREENSHOT_DIR / "validation-example.png",
            "Validation example",
        ),
        SlideSpec(
            6,
            "Access and run options",
            ["Run the app locally with one command from scripts/run-taskbuddy", "Use Docker or Docker Compose when preferred", "Build the documentation pack separately with its own .review-pack-venv"],
            notes["Access and run options"],
            30,
            SCREENSHOT_DIR / "login-page.png",
            "Sign-in screen",
        ),
        SlideSpec(
            7,
            "Architecture and LangGraph orchestration",
            ["FastAPI serves both the API and the built frontend", "AgentController keeps a stable contract for the route layer", "LangGraph drives validation, planning, tool execution, retry, and response assembly"],
            notes["Architecture and LangGraph orchestration"],
            35,
            DIAGRAM_DIR / "architecture-overview.png",
            "Architecture overview",
        ),
        SlideSpec(
            8,
            "Folder structure, key files, and APIs",
            ["backend/, frontend/, docs/, scripts/, tests/, and reports each have a focused role", "App entrypoints stay small while App.tsx, the interpreter, and the repository carry the main product logic", "The technical documentation lists every major API endpoint and its purpose"],
            notes["Folder structure, key files, and APIs"],
            45,
            DIAGRAM_DIR / "project-layout.png",
            "Project layout overview",
        ),
        SlideSpec(
            9,
            "Testing evidence and outputs",
            ["Backend and frontend tests generate JUnit XML outputs", "The report exporter builds an HTML dashboard and Excel summary", "The manual test plan and user guide complement automated checks"],
            notes["Testing evidence and outputs"],
            25,
            DIAGRAM_DIR / "test-evidence.png",
            "Testing evidence summary",
        ),
        SlideSpec(
            10,
            "Limits and next improvements",
            ["Current limits are documented in the README and technical guide", "The next technical improvements are model refinement, richer semantics, and more end-to-end automation", "The documentation pack is reproducible without touching the runtime environment"],
            notes["Limits and next improvements"],
            20,
            DIAGRAM_DIR / "request-lifecycle.png",
            "Request lifecycle",
        ),
    ]


def add_slide_notes(slide, notes: str) -> None:
    with contextlib.suppress(Exception):
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = notes


def add_full_background(slide, color: tuple[int, int, int]) -> None:
    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        PptInches(13.333),
        PptInches(7.5),
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PptColor(*color)
    background.line.fill.background()


def add_text_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int,
    color: tuple[int, int, int],
    bold: bool = False,
    align: int = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.size = PptPt(size)
    paragraph.font.name = "Segoe UI"
    paragraph.font.bold = bold
    paragraph.font.color.rgb = PptColor(*color)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_bullet_card(slide, left: float, top: float, width: float, height: float, bullets: list[str], fill: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptColor(*fill)
    shape.line.color.rgb = PptColor(*NAVY)

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = PptInches(0.18)
    frame.margin_top = PptInches(0.12)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(18)
        paragraph.font.name = "Segoe UI"
        paragraph.font.color.rgb = PptColor(*SLATE)
        paragraph.bullet = True
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = PptPt(10)


def add_image_card(slide, left: float, top: float, width: float, height: float, image_path: Path | None, caption: str | None) -> None:
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = PptColor(*WHITE)
    frame.line.color.rgb = PptColor(*NAVY)

    if image_path and image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            PptInches(left + 0.08),
            PptInches(top + 0.08),
            width=PptInches(width - 0.16),
            height=PptInches(height - 0.52),
        )
    else:
        add_text_box(slide, left + 0.2, top + 1.5, width - 0.4, 0.6, "Image unavailable", 20, SLATE, False, PP_ALIGN.CENTER)

    if caption:
        add_text_box(slide, left + 0.12, top + height - 0.42, width - 0.24, 0.3, caption, 9, SLATE, False, PP_ALIGN.CENTER)


def add_footer(slide, number: int) -> None:
    ribbon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(0),
        PptInches(7.12),
        PptInches(13.333),
        PptInches(0.38),
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = PptColor(*NAVY)
    ribbon.line.fill.background()
    add_text_box(slide, 0.48, 7.13, 4.0, 0.2, "TaskBuddy documentation pack", 9, WHITE, False)
    add_text_box(slide, 12.35, 7.13, 0.5, 0.2, str(number), 9, WHITE, True, PP_ALIGN.RIGHT)


def build_title_slide(slide, spec: SlideSpec) -> None:
    add_full_background(slide, NAVY)
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, PptInches(8.7), PptInches(-0.8), PptInches(4.4), PptInches(4.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PptColor(*TEAL)
    circle.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(0.7), PptInches(1.15), PptInches(0.9), PptInches(0.22))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PptColor(*ACCENT)
    accent.line.fill.background()

    add_text_box(slide, 0.8, 1.55, 6.0, 0.9, spec.title, 28, WHITE, True)
    add_text_box(slide, 0.82, 2.28, 5.8, 0.45, "Product walkthrough, technical design, and documentation pack", 15, (226, 235, 244))
    add_text_box(slide, 0.82, 2.82, 5.0, 0.35, spec.subtitle or "", 12, (226, 235, 244), False)
    add_bullet_card(slide, 0.82, 3.4, 5.25, 2.15, spec.bullets, TEAL_SOFT)
    add_image_card(slide, 7.15, 1.35, 5.45, 4.8, spec.image, spec.image_caption)
    add_text_box(slide, 0.82, 6.25, 7.0, 0.35, "Generated from the repo sources and the local FastAPI-served UI.", 10, (226, 235, 244))
    add_footer(slide, spec.number)


def build_content_slide(slide, spec: SlideSpec) -> None:
    add_full_background(slide, LIGHT_BG)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(13.333), PptInches(0.58))
    banner.fill.solid()
    banner.fill.fore_color.rgb = PptColor(*NAVY)
    banner.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, PptInches(11.85), PptInches(0.8), PptInches(0.45), PptInches(0.45))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PptColor(*ACCENT)
    accent.line.fill.background()

    add_text_box(slide, 0.68, 0.78, 7.0, 0.55, spec.title, 24, NAVY, True)
    add_text_box(slide, 0.7, 1.24, 1.3, 0.28, f"Target {int(spec.duration_seconds)}s", 10, TEAL, True)
    add_bullet_card(slide, 0.68, 1.7, 5.45, 4.75, spec.bullets, WHITE)
    add_image_card(slide, 6.48, 1.52, 6.12, 4.98, spec.image, spec.image_caption)
    add_footer(slide, spec.number)


def build_pptx() -> None:
    presentation = Presentation()
    presentation.slide_width = PptInches(13.333)
    presentation.slide_height = PptInches(7.5)
    blank_layout = presentation.slide_layouts[6]

    for spec in slide_specs():
        slide = presentation.slides.add_slide(blank_layout)
        if spec.number == 1:
            build_title_slide(slide, spec)
        else:
            build_content_slide(slide, spec)
        add_slide_notes(slide, spec.notes)

    presentation.save(OUTPUT_DIR / "TaskBuddy-Demo-Deck.pptx")


def create_video_slide_image(spec: SlideSpec) -> Path:
    width, height = 1600, 900
    image = Image.new("RGB", (width, height), LIGHT_BG)
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((40, 40, width - 40, height - 40), radius=34, fill=WHITE, outline=NAVY, width=4)
    draw.rounded_rectangle((40, 40, width - 40, 130), radius=34, fill=NAVY, outline=NAVY, width=0)
    draw.text((88, 72), f"Slide {spec.number} • {spec.title}", font=load_font(36, bold=True), fill=WHITE)
    draw.text((88, 155), "TaskBuddy five-minute walkthrough", font=load_font(20), fill=SLATE)

    draw.rounded_rectangle((88, 210, 700, 690), radius=24, fill=NAVY_SOFT, outline=NAVY, width=3)
    y = 252
    bullet_font = load_font(28)
    for bullet in spec.bullets:
        wrapped = textwrap.wrap(bullet, width=34)
        for line_index, line in enumerate(wrapped):
            prefix = "- " if line_index == 0 else "  "
            draw.text((118, y), prefix + line, font=bullet_font, fill=SLATE)
            y += 38
        y += 10

    if spec.image and spec.image.exists():
        preview = Image.open(spec.image).convert("RGB")
        preview.thumbnail((680, 460))
        px = width - preview.width - 100
        py = 210
        draw.rounded_rectangle((px - 10, py - 10, px + preview.width + 10, py + preview.height + 10), radius=22, fill=WHITE, outline=NAVY, width=3)
        image.paste(preview, (px, py))
        if spec.image_caption:
            draw.text((px, py + preview.height + 28), spec.image_caption, font=load_font(20), fill=SLATE)

    draw.text((88, 760), f"Narration target: {int(spec.duration_seconds)} seconds", font=load_font(18), fill=TEAL)
    output_path = VIDEO_DIR / f"slide-{spec.number:02d}.png"
    image.save(output_path)
    return output_path


def create_video_cards(slides: list[SlideSpec]) -> list[Path]:
    return [create_video_slide_image(spec) for spec in slides]


def wait_for_audio_files(paths: Iterable[Path], timeout_seconds: float = 5.0) -> None:
    deadline = time.time() + timeout_seconds
    while time.time() < deadline:
        if all(path.exists() and path.stat().st_size > 0 for path in paths):
            return
        time.sleep(0.2)
    raise RuntimeError("Audio files were not written before the timeout expired.")


def generate_pyttsx3_audio(slides: list[SlideSpec]) -> list[Path]:
    import pyttsx3

    engine = pyttsx3.init()
    engine.setProperty("rate", 185)
    output_paths: list[Path] = []
    try:
        for spec in slides:
            output_path = VIDEO_DIR / f"slide-{spec.number:02d}.wav"
            if output_path.exists():
                output_path.unlink()
            engine.save_to_file(spec.notes, str(output_path))
            output_paths.append(output_path)
        engine.runAndWait()
        wait_for_audio_files(output_paths)
    finally:
        with contextlib.suppress(Exception):
            engine.stop()
    return output_paths


def generate_gtts_audio(slides: list[SlideSpec]) -> list[Path]:
    from gtts import gTTS

    output_paths: list[Path] = []
    for spec in slides:
        output_path = VIDEO_DIR / f"slide-{spec.number:02d}.mp3"
        if output_path.exists():
            output_path.unlink()
        audio = gTTS(text=spec.notes, lang="en", slow=False)
        audio.save(str(output_path))
        output_paths.append(output_path)
    wait_for_audio_files(output_paths)
    return output_paths


def generate_voiceover_files(slides: list[SlideSpec]) -> tuple[str | None, list[Path], list[str]]:
    status_lines: list[str] = []
    for existing in VIDEO_DIR.glob("slide-*.wav"):
        existing.unlink()
    for existing in VIDEO_DIR.glob("slide-*.mp3"):
        existing.unlink()

    try:
        return "pyttsx3", generate_pyttsx3_audio(slides), status_lines
    except Exception as error:
        status_lines.append(f"Offline narration failed with pyttsx3: {error}")

    try:
        return "gTTS", generate_gtts_audio(slides), status_lines
    except Exception as error:
        status_lines.append(f"Online narration fallback failed with gTTS: {error}")

    return None, [], status_lines


def format_duration(seconds: float) -> str:
    total_seconds = int(round(seconds))
    minutes, remaining = divmod(total_seconds, 60)
    return f"{minutes:02d}:{remaining:02d}"


def attempt_video(slides: list[SlideSpec]) -> None:
    status_lines: list[str] = []
    video_path = OUTPUT_DIR / "TaskBuddy-Demo.mp4"
    provider: str | None = None
    audio_paths: list[Path] = []
    card_paths = create_video_cards(slides)
    silent_fallback = False
    total_duration = 0.0
    clips = []
    audio_clips = []

    try:
        provider, audio_paths, provider_status = generate_voiceover_files(slides)
        status_lines.extend(provider_status)

        try:
            from moviepy import AudioFileClip, ImageClip, concatenate_videoclips
        except ImportError:
            from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips

        if provider and len(audio_paths) == len(card_paths):
            for card_path, audio_path in zip(card_paths, audio_paths, strict=True):
                audio_clip = AudioFileClip(str(audio_path))
                audio_clips.append(audio_clip)
                clip = ImageClip(str(card_path))
                duration = audio_clip.duration + 0.2
                total_duration += duration
                try:
                    clip = clip.with_duration(duration).with_audio(audio_clip)
                except AttributeError:
                    clip = clip.set_duration(duration).set_audio(audio_clip)
                clips.append(clip)
            status_lines.append(f"Narration provider: {provider}")
        else:
            silent_fallback = True
            status_lines.append("Audio narration was unavailable. Rendered a silent slideshow fallback using the target slide durations.")
            for spec, card_path in zip(slides, card_paths, strict=True):
                total_duration += spec.duration_seconds
                clip = ImageClip(str(card_path))
                try:
                    clip = clip.with_duration(spec.duration_seconds)
                except AttributeError:
                    clip = clip.set_duration(spec.duration_seconds)
                clips.append(clip)

        final_clip = concatenate_videoclips(clips, method="compose")
        final_clip.write_videofile(
            str(video_path),
            fps=24,
            codec="libx264",
            audio_codec="aac",
            logger=None,
        )
        final_clip.close()

        status_lines.append(f"Total demo duration: {format_duration(total_duration)}")
        if not silent_fallback and not 270 <= total_duration <= 300:
            status_lines.append("Narrated video length is outside the 04:30 to 05:00 target window.")
        if silent_fallback:
            status_lines.append("The PPTX deck and DOCX demo script remain the primary narrated presentation assets.")
    except Exception as error:
        status_lines.append(f"MP4 generation did not complete cleanly: {error}")
    finally:
        for clip in clips:
            with contextlib.suppress(Exception):
                clip.close()
        for audio_clip in audio_clips:
            with contextlib.suppress(Exception):
                audio_clip.close()

    (OUTPUT_DIR / "TaskBuddy-Demo-Video-Status.txt").write_text("\n".join(status_lines) + "\n", encoding="utf-8")


def build_demo_script_docx() -> None:
    generate_docx(
        DOCS_DIR / "demo-script.md",
        OUTPUT_DIR / "TaskBuddy-Demo-Script.docx",
        "Demo Script",
        "Narration notes for the five-minute product and technical walkthrough.",
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate the TaskBuddy documentation pack.")
    parser.add_argument("--app-url", default="http://localhost:8000", help="Base URL of a running TaskBuddy app.")
    parser.add_argument("--skip-video", action="store_true", help="Skip MP4 generation.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    base_url = args.app_url.rstrip("/")
    ensure_output_dirs()
    generate_diagrams()
    capture_screenshots(base_url)

    generate_docx(
        DOCS_DIR / "user-guide.md",
        OUTPUT_DIR / "TaskBuddy-User-Guide.docx",
        "User Guide",
        "Setup instructions, core workflows, response interpretation, and troubleshooting guidance.",
    )
    generate_docx(
        DOCS_DIR / "technical-design.md",
        OUTPUT_DIR / "TaskBuddy-Technical-Documentation.docx",
        "Technical Documentation",
        "Architecture, folder structure, API purposes, orchestration, persistence, and testing details.",
    )
    generate_docx(
        DOCS_DIR / "manual-test-plan.md",
        OUTPUT_DIR / "TaskBuddy-Manual-Test-Plan.docx",
        "Manual Test Plan",
        "Structured tool checks and workflow validations for local verification.",
    )
    build_demo_script_docx()
    build_pptx()
    if not args.skip_video:
        attempt_video(slide_specs())


if __name__ == "__main__":
    main()
